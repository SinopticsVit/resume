"""
Сборка PPTX-презентации «WFOE Детский Мир в КНР — Фаза 1 (6 месяцев)».

Выход: DetMir_China_WFOE_Phase1.pptx (15 слайдов, 16:9, фирменная палитра).

Mermaid-диаграммы берутся из _assets/diagrams/*.mmd и рендерятся в PNG:
  1) npx @mermaid-js/mermaid-cli (если установлен Node.js)
  2) иначе — mermaid.ink (HTTP, без Node.js)
полный контроль вёрстки, без хрупкого парсинга markdown.

Использование:
    python -m pip install -r requirements.txt
    python build_pptx.py                 # render mermaid + build pptx
    python build_pptx.py --skip-mermaid  # reuse cached PNGs (без Node.js)

Спецификация раскладки: build-spec.md; источник текста: slides.md.
"""

from __future__ import annotations

import argparse
import base64
import os
import shutil
import subprocess
import sys
import time
import urllib.error
import urllib.request
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# Windows-консоль (cp1252) не печатает кириллицу — переключаем потоки на UTF-8.
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except (AttributeError, ValueError):
        pass

# ---------------------------------------------------------------------------
# Пути
# ---------------------------------------------------------------------------

HERE = Path(__file__).resolve().parent
ASSETS_DIR = HERE / "_assets" / "diagrams"
OUTPUT_PPTX = HERE / "DetMir_China_WFOE_Phase1.pptx"

# Файлы диаграмм в порядке использования в деке.
DIAGRAM_FILES = [
    "diagram-01-func-blocks.mmd",
    "diagram-02-org-structure.mmd",
    "diagram-03-budget-flow.mmd",
    "diagram-04-roadmap-phases.mmd",
    "diagram-05-moscow-links.mmd",
    "diagram-06-process-samples.mmd",
    "diagram-07-process-qc.mmd",
    "diagram-08-sample-payment.mmd",
]

# ---------------------------------------------------------------------------
# Дизайн-токены (палитра ритейл/закупки: navy + teal + crimson + gold + ivory)
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_DARK = RGBColor(0x12, 0x1D, 0x36)
NAVY_LIGHT = RGBColor(0x2C, 0x3E, 0x6B)
TEAL = RGBColor(0x1A, 0x7A, 0x6D)
TEAL_SOFT = RGBColor(0xEA, 0xF5, 0xF3)
CRIMSON = RGBColor(0xB0, 0x2A, 0x30)
CRIMSON_SOFT = RGBColor(0xFB, 0xEE, 0xEE)
GOLD = RGBColor(0xC8, 0x96, 0x2E)
GOLD_SOFT = RGBColor(0xFB, 0xF3, 0xDF)

BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_IVORY = RGBColor(0xFA, 0xF8, 0xF5)
BG_SOFT = RGBColor(0xF2, 0xEF, 0xEB)
RULE = RGBColor(0xCC, 0xC5, 0xB9)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_BODY = RGBColor(0x33, 0x33, 0x33)
TEXT_MUTED = RGBColor(0x66, 0x66, 0x66)

TABLE_HEADER_BG = NAVY
TABLE_ROW_EVEN = BG_SOFT
TABLE_ROW_ODD = BG_IVORY

FONT_HEAD = "Georgia"
FONT_BODY = "Segoe UI"

HEADER_H = Inches(1.15)
FOOTER_H = Inches(0.32)
SIDE_M = Inches(0.7)
CONTENT_TOP = HEADER_H + Inches(0.35)
DECK_LABEL = "Детский Мир · WFOE КНР · Фаза 1"
DOC_TITLE = "Сервисный центр закупок в Китае · 6 месяцев"

# ---------------------------------------------------------------------------
# Рендеринг mermaid
# ---------------------------------------------------------------------------

def _resolve_npx() -> Optional[str]:
    candidates = ["npx.cmd", "npx.exe", "npx"] if os.name == "nt" else ["npx"]
    for name in candidates:
        path = shutil.which(name)
        if path:
            return path
    return None


def _render_mermaid_cli(npx: str, mmd: Path, out: Path) -> None:
    cmd = [
        npx, "-y", "@mermaid-js/mermaid-cli",
        "-i", str(mmd),
        "-o", str(out),
        "-b", "transparent",
        "-w", "2000",
        "-s", "2",
    ]
    use_shell = os.name == "nt" and npx.lower().endswith(".cmd")
    proc = subprocess.run(cmd, capture_output=True, text=True, shell=use_shell)
    if proc.returncode != 0:
        print(proc.stdout)
        print(proc.stderr, file=sys.stderr)
        raise RuntimeError(f"mermaid-cli failed for {out.name}")


def _render_mermaid_ink(mmd: Path, out: Path) -> None:
    """Fallback: рендер через публичный сервис mermaid.ink (без Node.js)."""
    src = mmd.read_text(encoding="utf-8")
    encoded = base64.urlsafe_b64encode(src.encode("utf-8")).decode("ascii").rstrip("=")
    url = f"https://mermaid.ink/img/{encoded}"
    req = urllib.request.Request(url, headers={"User-Agent": "DetMir-build-pptx/1.0"})
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = resp.read()
    except urllib.error.URLError as exc:
        raise RuntimeError(f"mermaid.ink недоступен для {mmd.name}: {exc}") from exc
    if len(data) < 500:
        raise RuntimeError(f"mermaid.ink вернул пустой ответ для {mmd.name}")
    out.write_bytes(data)


def render_mermaid(*, skip: bool) -> list[Path]:
    """Рендер каждого .mmd в PNG. PNG лежит рядом с .mmd с суффиксом .png."""
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    npx = _resolve_npx()
    paths: list[Path] = []
    for fname in DIAGRAM_FILES:
        mmd = ASSETS_DIR / fname
        out = mmd.with_suffix(".png")
        paths.append(out)
        if not mmd.exists():
            print(f"[warn] нет исходника диаграммы: {mmd.name}", file=sys.stderr)
            continue
        if out.exists() and skip:
            print(f"[cache] {out.name}")
            continue
        if npx is not None:
            print(f"[mermaid-cli] {mmd.name} -> {out.name}")
            _render_mermaid_cli(npx, mmd, out)
        else:
            print(f"[mermaid.ink] {mmd.name} -> {out.name}")
            _render_mermaid_ink(mmd, out)
    return paths


# ---------------------------------------------------------------------------
# Низкоуровневые helpers
# ---------------------------------------------------------------------------

def _solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def _set_text(
    tf,
    text: str,
    *,
    font: str = FONT_BODY,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_BODY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    spacing: float = 1.1,
) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _new_blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid_fill(bg, BG_IVORY)
    return slide


def _add_header(slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    _solid_fill(band, NAVY)
    label = slide.shapes.add_textbox(SIDE_M, Inches(0.16), Inches(9.0), Inches(0.3))
    _set_text(label.text_frame, DECK_LABEL, font=FONT_BODY, size=11,
              color=RGBColor(0xB9, 0xC2, 0xD6))
    title_box = slide.shapes.add_textbox(
        SIDE_M, Inches(0.48), SLIDE_W - 2 * SIDE_M, Inches(0.6)
    )
    _set_text(title_box.text_frame, title, font=FONT_HEAD, size=26,
              bold=True, color=BG, anchor=MSO_ANCHOR.MIDDLE)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, SIDE_M, Inches(1.0), Inches(2.6), Pt(4)
    )
    _solid_fill(accent, GOLD)


def _add_footer(slide, page_num: int, page_total: int) -> None:
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - FOOTER_H, SLIDE_W, Emu(1)
    )
    _solid_fill(rule, RULE)
    left = slide.shapes.add_textbox(
        SIDE_M, SLIDE_H - FOOTER_H + Inches(0.04), Inches(9.0), Inches(0.25)
    )
    _set_text(left.text_frame, DOC_TITLE, font=FONT_BODY, size=10, color=TEXT_MUTED)
    right = slide.shapes.add_textbox(
        SLIDE_W - Inches(2.0), SLIDE_H - FOOTER_H + Inches(0.04),
        Inches(1.3), Inches(0.25),
    )
    _set_text(right.text_frame, f"{page_num} / {page_total}",
              font=FONT_BODY, size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def _content_area():
    x = SIDE_M
    y = CONTENT_TOP
    w = SLIDE_W - 2 * SIDE_M
    h = SLIDE_H - CONTENT_TOP - FOOTER_H - Inches(0.1)
    return x, y, w, h


def _bullets(slide, x, y, w, items, *, size=15, gap=0.5, lead_color=NAVY,
             body_color=TEXT_BODY, bullet_char="•", bullet_color=CRIMSON):
    """items: список строк или кортежей (lead, rest)."""
    cur_y = y
    for item in items:
        line_h = Inches(gap)
        box = slide.shapes.add_textbox(x, cur_y, w, line_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.text = ""
        p.line_spacing = 1.05
        b = p.add_run(); b.text = f"{bullet_char}  "
        b.font.name = FONT_BODY; b.font.size = Pt(size); b.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            lead, rest = item
            if lead:
                lr = p.add_run(); lr.text = lead
                lr.font.name = FONT_BODY; lr.font.size = Pt(size)
                lr.font.bold = True; lr.font.color.rgb = lead_color
            rr = p.add_run(); rr.text = rest
            rr.font.name = FONT_BODY; rr.font.size = Pt(size)
            rr.font.color.rgb = body_color
        else:
            rr = p.add_run(); rr.text = item
            rr.font.name = FONT_BODY; rr.font.size = Pt(size)
            rr.font.color.rgb = body_color
        cur_y += line_h
    return cur_y


def _accent_box(slide, x, y, w, text, *, fill=TEAL_SOFT, border=TEAL,
                text_color=TEAL, h=0.9, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border; box.line.width = Pt(1)
    box.shadow.inherit = False
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.12),
                                  w - Inches(0.5), Inches(h) - Inches(0.24))
    _set_text(tb.text_frame, text, font=FONT_BODY, size=size, bold=True,
              color=text_color, anchor=MSO_ANCHOR.MIDDLE)


def _embed_image(slide, image_path: Path, x, y, w, h, *, framed=True):
    if framed:
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        frame.fill.solid(); frame.fill.fore_color.rgb = BG
        frame.line.color.rgb = RULE; frame.line.width = Pt(1)
        frame.shadow.inherit = False
        frame.adjustments[0] = 0.03
    if image_path.exists():
        try:
            from PIL import Image
            with Image.open(image_path) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = 16, 9
        pad = Inches(0.25)
        max_w = w - 2 * pad
        max_h = h - 2 * pad
        ratio = min(max_w / iw, max_h / ih)
        tw = int(iw * ratio); th = int(ih * ratio)
        px = x + (w - tw) / 2; py = y + (h - th) / 2
        slide.shapes.add_picture(str(image_path), px, py, tw, th)
    else:
        ph = slide.shapes.add_textbox(x, y, w, h)
        _set_text(ph.text_frame, f"[нет диаграммы: {image_path.name}]",
                  font=FONT_BODY, size=14, color=CRIMSON,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _native_table(slide, x, y, w, headers, rows, col_widths, *,
                  max_h=None, header_size=12, body_size=11,
                  right_align_cols=None, emphasis_rows=None):
    right_align_cols = right_align_cols or set()
    emphasis_rows = emphasis_rows or set()
    nrows = len(rows) + 1
    ncols = len(headers)
    row_h = Inches(0.42)
    table_h = row_h * nrows
    if max_h is not None and table_h > max_h:
        table_h = max_h
    shape = slide.shapes.add_table(nrows, ncols, x, y, w, table_h)
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    total = sum(col_widths)
    for i, cw in enumerate(col_widths):
        table.columns[i].width = Emu(int(w * (cw / total)))
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HEADER_BG
        cell.margin_top = Emu(0); cell.margin_bottom = Emu(0)
        align = PP_ALIGN.RIGHT if ci in right_align_cols else PP_ALIGN.LEFT
        _set_text(cell.text_frame, header, font=FONT_HEAD, size=header_size,
                  bold=True, color=BG, align=align, anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        emphasized = (ri - 1) in emphasis_rows
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            if emphasized:
                cell.fill.fore_color.rgb = GOLD_SOFT
            else:
                cell.fill.fore_color.rgb = TABLE_ROW_EVEN if ri % 2 == 0 else TABLE_ROW_ODD
            cell.margin_top = Emu(0); cell.margin_bottom = Emu(0)
            align = PP_ALIGN.RIGHT if ci in right_align_cols else PP_ALIGN.LEFT
            _set_text(cell.text_frame, val, font=FONT_BODY, size=body_size,
                      color=TEXT_BODY, bold=emphasized, align=align,
                      anchor=MSO_ANCHOR.MIDDLE)
    return shape


# ---------------------------------------------------------------------------
# Сборщики слайдов
# ---------------------------------------------------------------------------

def slide_title(prs, *, eyebrow, title, subtitle, meta, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid_fill(bg, NAVY_DARK)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.16))
    _solid_fill(top, GOLD)
    crimson = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.16), SLIDE_W, Inches(0.06))
    _solid_fill(crimson, CRIMSON)

    eb = slide.shapes.add_textbox(SIDE_M, Inches(1.1), SLIDE_W - 2 * SIDE_M, Inches(0.5))
    _set_text(eb.text_frame, eyebrow, font=FONT_BODY, size=15,
              color=RGBColor(0xC9, 0xB8, 0x7A))
    tb = slide.shapes.add_textbox(SIDE_M, Inches(1.9), SLIDE_W - 2 * SIDE_M, Inches(2.2))
    _set_text(tb.text_frame, title, font=FONT_HEAD, size=40, bold=True,
              color=BG, spacing=1.15)
    _solid_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SIDE_M, Inches(4.3),
                               Inches(3.5), Pt(4)),
        CRIMSON,
    )
    sb = slide.shapes.add_textbox(SIDE_M, Inches(4.6), SLIDE_W - 2 * SIDE_M, Inches(1.2))
    _set_text(sb.text_frame, subtitle, font=FONT_BODY, size=17,
              color=RGBColor(0xC9, 0xCF, 0xDC), spacing=1.3)
    mb = slide.shapes.add_textbox(SIDE_M, Inches(5.9), SLIDE_W - 2 * SIDE_M, Inches(0.5))
    _set_text(mb.text_frame, meta, font=FONT_BODY, size=14,
              color=RGBColor(0xA7, 0xB0, 0xC4))

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.7),
                                   SLIDE_W, Inches(0.7))
    _solid_fill(strip, NAVY)
    fb = slide.shapes.add_textbox(SIDE_M, SLIDE_H - Inches(0.6),
                                  SLIDE_W - 2 * SIDE_M, Inches(0.45))
    _set_text(fb.text_frame, footer, font=FONT_BODY, size=13,
              color=RGBColor(0xB9, 0xC2, 0xD6), align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)


def slide_header_bullets(prs, *, page, total, title, intro=None, bullets,
                         accent=None, accent_kind="teal", size=15, gap=0.55):
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    cur_y = y
    if intro:
        box = slide.shapes.add_textbox(x, cur_y, w, Inches(0.85))
        _set_text(box.text_frame, intro, font=FONT_BODY, size=15,
                  color=TEXT_DARK, spacing=1.25)
        cur_y += Inches(1.0)
    cur_y = _bullets(slide, x, cur_y, w, bullets, size=size, gap=gap)
    if accent:
        palette = {
            "teal": (TEAL_SOFT, TEAL, TEAL),
            "crimson": (CRIMSON_SOFT, CRIMSON, CRIMSON),
            "gold": (GOLD_SOFT, GOLD, NAVY),
        }
        fill, border, tcolor = palette.get(accent_kind, palette["teal"])
        _accent_box(slide, x, cur_y + Inches(0.2), w, accent,
                    fill=fill, border=border, text_color=tcolor)


def slide_diagram(prs, *, page, total, title, image_path, intro=None,
                  caption=None, bullets=None):
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    cur_y = y
    if intro:
        box = slide.shapes.add_textbox(x, cur_y, w, Inches(0.55))
        _set_text(box.text_frame, intro, font=FONT_BODY, size=14,
                  color=TEXT_MUTED, spacing=1.2)
        cur_y += Inches(0.6)
    bullets_h = Inches(0.0)
    if bullets:
        bullets_h = Inches(0.42 * len(bullets) + 0.2)
    cap_h = Inches(0.4) if caption else Inches(0.0)
    img_h = h - (cur_y - y) - cap_h - bullets_h - Inches(0.1)
    _embed_image(slide, image_path, x, cur_y, w, img_h)
    cur_y += img_h
    if caption:
        cb = slide.shapes.add_textbox(x, cur_y + Inches(0.03), w, cap_h)
        _set_text(cb.text_frame, caption, font=FONT_BODY, size=12,
                  color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        cur_y += cap_h
    if bullets:
        _bullets(slide, x, cur_y + Inches(0.05), w, bullets, size=13, gap=0.42)


def slide_diagram_right(prs, *, page, total, title, image_path, bullets,
                        intro=None):
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    left_w = w * 0.42
    right_x = x + left_w + Inches(0.3)
    right_w = w - left_w - Inches(0.3)
    cur_y = y
    if intro:
        box = slide.shapes.add_textbox(x, cur_y, left_w, Inches(0.6))
        _set_text(box.text_frame, intro, font=FONT_BODY, size=14,
                  color=TEXT_MUTED, spacing=1.2)
        cur_y += Inches(0.7)
    _bullets(slide, x, cur_y, left_w, bullets, size=14, gap=0.7)
    _embed_image(slide, image_path, right_x, y, right_w, h)


def slide_table(prs, *, page, total, title, headers, rows, col_widths,
                intro=None, note=None, right_align_cols=None,
                emphasis_rows=None, body_size=11):
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    cur_y = y
    if intro:
        box = slide.shapes.add_textbox(x, cur_y, w, Inches(0.55))
        _set_text(box.text_frame, intro, font=FONT_BODY, size=14,
                  color=TEXT_MUTED, spacing=1.2)
        cur_y += Inches(0.6)
    note_h = Inches(0.55) if note else Inches(0.0)
    _native_table(slide, x, cur_y, w, headers, rows, col_widths,
                  max_h=h - (cur_y - y) - note_h,
                  right_align_cols=right_align_cols,
                  emphasis_rows=emphasis_rows, body_size=body_size)
    if note:
        nb = slide.shapes.add_textbox(x, SLIDE_H - FOOTER_H - Inches(0.55),
                                      w, Inches(0.5))
        _set_text(nb.text_frame, note, font=FONT_BODY, size=11,
                  color=TEAL, spacing=1.15)


def slide_two_column_tables(prs, *, page, total, title, left_title, left_headers,
                            left_rows, left_cw, right_title, right_headers,
                            right_rows, right_cw, note=None):
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    col_w = (w - Inches(0.4)) / 2
    right_x = x + col_w + Inches(0.4)
    note_h = Inches(0.5) if note else Inches(0.0)
    for cx, ctitle, headers, rows, cw, color in [
        (x, left_title, left_headers, left_rows, left_cw, TEAL),
        (right_x, right_title, right_headers, right_rows, right_cw, NAVY_LIGHT),
    ]:
        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, col_w, Inches(0.45))
        _solid_fill(lbl, color)
        _set_text(lbl.text_frame, ctitle, font=FONT_BODY, size=14, bold=True,
                  color=BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _native_table(slide, cx, y + Inches(0.6), col_w, headers, rows, cw,
                      max_h=h - Inches(0.6) - note_h, body_size=10.5,
                      header_size=11)
    if note:
        nb = slide.shapes.add_textbox(x, SLIDE_H - FOOTER_H - Inches(0.5),
                                      w, Inches(0.45))
        _set_text(nb.text_frame, note, font=FONT_BODY, size=11,
                  color=TEAL, spacing=1.1)


def slide_diagram_table(prs, *, page, total, title, image_path, headers, rows,
                        col_widths, intro=None, emphasis_rows=None):
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    cur_y = y
    if intro:
        box = slide.shapes.add_textbox(x, cur_y, w, Inches(0.5))
        _set_text(box.text_frame, intro, font=FONT_BODY, size=13,
                  color=TEXT_MUTED)
        cur_y += Inches(0.55)
    img_h = Inches(1.7)
    _embed_image(slide, image_path, x, cur_y, w, img_h)
    cur_y += img_h + Inches(0.2)
    _native_table(slide, x, cur_y, w, headers, rows, col_widths,
                  max_h=h - (cur_y - y), body_size=10.5,
                  emphasis_rows=emphasis_rows)


def slide_three_processes(prs, *, page, total, title, images, captions, note=None):
    """Слайд 14: два процесса сверху (две колонки) + контур оплаты снизу."""
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    note_h = Inches(0.5) if note else Inches(0.0)
    body_h = h - note_h
    top_h = body_h * 0.5
    bottom_h = body_h * 0.5 - Inches(0.1)
    col_w = (w - Inches(0.3)) / 2
    right_x = x + col_w + Inches(0.3)
    # верхний ряд — два процесса
    for cx, img, cap in [
        (x, images[0], captions[0]),
        (right_x, images[1], captions[1]),
    ]:
        lbl = slide.shapes.add_textbox(cx, y, col_w, Inches(0.3))
        _set_text(lbl.text_frame, cap, font=FONT_BODY, size=13, bold=True,
                  color=NAVY)
        _embed_image(slide, img, cx, y + Inches(0.32), col_w, top_h - Inches(0.32))
    # нижний ряд — контур оплаты
    by = y + top_h + Inches(0.1)
    lbl = slide.shapes.add_textbox(x, by, w, Inches(0.3))
    _set_text(lbl.text_frame, captions[2], font=FONT_BODY, size=13, bold=True,
              color=NAVY)
    _embed_image(slide, images[2], x, by + Inches(0.32), w, bottom_h - Inches(0.32))
    if note:
        nb = slide.shapes.add_textbox(x, SLIDE_H - FOOTER_H - Inches(0.5),
                                      w, Inches(0.45))
        _set_text(nb.text_frame, note, font=FONT_BODY, size=11,
                  color=TEAL, spacing=1.1)


def slide_closing(prs, *, page, total, title, headline, sla_rows, risk_rows,
                  checklist):
    """Слайд 15: SLA + риски + критерии Фазы 2 + финальный тезис."""
    slide = _new_blank(prs)
    _add_header(slide, title)
    _add_footer(slide, page, total)
    x, y, w, h = _content_area()
    col_w = (w - Inches(0.6)) / 3
    mid_x = x + col_w + Inches(0.3)
    right_x = x + 2 * (col_w + Inches(0.3))
    table_top = y + Inches(0.4)
    table_h = Inches(3.4)
    # колонка 1 — SLA
    _col_caption(slide, x, y, col_w, "SLA (ключевые)", TEAL)
    _native_table(slide, x, table_top, col_w, ["Параметр", "Значение"],
                  sla_rows, [2.0, 1.0], max_h=table_h, body_size=9.5,
                  header_size=10)
    # колонка 2 — риски
    _col_caption(slide, mid_x, y, col_w, "Риски и митигация", CRIMSON)
    _native_table(slide, mid_x, table_top, col_w, ["Риск", "Митигация"],
                  risk_rows, [1.2, 1.4], max_h=table_h, body_size=9.5,
                  header_size=10)
    # колонка 3 — чеклист Фазы 2
    _col_caption(slide, right_x, y, col_w, "Критерии Фазы 2 (M6)", GOLD)
    cy = table_top
    for item in checklist:
        box = slide.shapes.add_textbox(right_x, cy, col_w, Inches(0.62))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = ""; p.line_spacing = 1.0
        chk = p.add_run(); chk.text = "☐  "
        chk.font.name = FONT_BODY; chk.font.size = Pt(11); chk.font.color.rgb = GOLD
        rr = p.add_run(); rr.text = item
        rr.font.name = FONT_BODY; rr.font.size = Pt(10.5); rr.font.color.rgb = TEXT_BODY
        cy += Inches(0.62)
    # финальный тезис
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,
                                 SLIDE_H - FOOTER_H - Inches(1.0), w, Inches(0.85))
    _solid_fill(box, NAVY)
    tb = slide.shapes.add_textbox(x + Inches(0.3), SLIDE_H - FOOTER_H - Inches(0.92),
                                  w - Inches(0.6), Inches(0.7))
    _set_text(tb.text_frame, headline, font=FONT_HEAD, size=14, bold=True,
              color=BG, anchor=MSO_ANCHOR.MIDDLE, spacing=1.15)


def _col_caption(slide, x, y, w, text, color):
    lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.38))
    _solid_fill(lbl, color)
    _set_text(lbl.text_frame, text, font=FONT_BODY, size=12, bold=True,
              color=BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Сохранение (устойчиво к открытому в PowerPoint файлу)
# ---------------------------------------------------------------------------

def _save_pptx(prs: Presentation, target: Path) -> Path:
    target.parent.mkdir(parents=True, exist_ok=True)
    tmp = target.with_suffix(target.suffix + ".tmp")
    if tmp.exists():
        try:
            tmp.unlink()
        except PermissionError:
            tmp = target.with_name(target.stem + f".{int(time.time())}.tmp")
    prs.save(str(tmp))
    for attempt in range(5):
        try:
            if target.exists():
                target.unlink()
            os.replace(tmp, target)
            return target
        except PermissionError:
            time.sleep(0.6 * (attempt + 1))
    fallback = target.with_name(
        target.stem + "." + time.strftime("%Y%m%d-%H%M%S") + target.suffix
    )
    os.replace(tmp, fallback)
    print(f"[warn] {target.name} заблокирован. Записан {fallback.name}.",
          file=sys.stderr)
    return fallback


# ---------------------------------------------------------------------------
# Композиция дека
# ---------------------------------------------------------------------------

def build_deck(diagrams: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    def dia(i: int) -> Path:
        return diagrams[i] if i < len(diagrams) else ASSETS_DIR / f"missing-{i}.png"

    TOTAL = 14  # слайды 2–15 нумеруются 1..14 в футере; титул без номера
    p = 1

    # --- Слайд 1. Титул ---
    slide_title(
        prs,
        eyebrow="Запуск компании в КНР · по итогам исследования office_setup",
        title="Сервисный центр закупок\nв Китае: WFOE «Детский Мир»",
        subtitle="Структура, бюджет и план запуска на 6 месяцев\nФаза 1 — услуги для московских ПМ, без торговли на собственный баланс",
        meta="WFOE (торговая компания с правом экспорта) · HQ Шанхай · QC-узел Гуанчжоу · [дата] · [докладчик]",
        footer="Сорсинг · образцы · ОТК · контент · тестирование · аналитика — по сервисным договорам (cost-plus)",
    )

    # --- Слайд 2. Миссия и модель Фазы 1 ---
    slide_header_bullets(
        prs, page=p, total=TOTAL,
        title="Миссия и модель Фазы 1: «глаза и руки» в Китае",
        intro="Миссия: дать московскому офису прямой, контролируемый и быстрый доступ к китайским фабрикам — сорсинг, образцы, ОТК, контент, тестирование, аналитика, разрешение проблем с поставщиками.",
        bullets=[
            ("Юридическая форма — ", "WFOE с правом экспорта, но 6 месяцев — только сервисные услуги по договорам с материнской компанией."),
            ("Самофинансирование — ", "выручка = cost base × 1,08 (cost-plus) + НДС 6%."),
            ("Русское управление + китайская команда — ", "директор и технолог (RU), операционный штат (CN)."),
            ("Экспортная лицензия — ", "оформляется «про запас»; активация в Фазе 2 (решение в конце M6)."),
        ],
        accent="Не представительство и не торговый дом «с первого дня», а контролируемый сервисный контур с прозрачным биллингом.",
        accent_kind="teal", size=15, gap=0.72,
    ); p += 1

    # --- Слайд 3. Scope in/out ---
    slide_two_column_tables(
        prs, page=p, total=TOTAL,
        title="Scope Фазы 1: что делаем и чего не делаем",
        left_title="В scope (6 месяцев)",
        left_headers=["Направление", "Содержание"],
        left_rows=[
            ["Sourcing-саппорт", "Поиск и верификация фабрик, AVL, КП по брифам ПМ"],
            ["Образцы", "Преселекция, консолидация, авиа в РФ (3–7 дней)"],
            ["ОТК и аудиты", "Pre-shipment AQL 2,5, аудиты фабрик/складов"],
            ["Брак и задержки", "Выезды на фабрику, CAPA, эскалация в Москву"],
            ["Контент/тесты/аналитика", "Фото/видео, лаборатории, тренды рынка"],
            ["Liaison + финансы", "Единое окно для ПМ, биллинг услуг"],
        ],
        left_cw=[1.2, 2.0],
        right_title="Out of scope (Фаза 1)",
        right_headers=["Исключено", "Комментарий"],
        right_rows=[
            ["Экспорт/торговля на баланс WFOE", "Лицензия есть, не используется"],
            ["Контракты на закупку/импорт", "Ведёт головная компания РФ"],
            ["Прямые платежи фабрикам за товар", "Исключение — pass-through лабораторий"],
            ["Складирование и таможня партий", "Функция логистики РФ"],
            ["B2C / маркетплейсы в КНР", "—"],
        ],
        right_cw=[1.4, 1.4],
        note="Фаза 2 (конец M6): go/no-go по активации торговых/экспортных операций без смены юр. формы.",
    ); p += 1

    # --- Слайд 4. Карта функций F1–F10 ---
    slide_diagram_right(
        prs, page=p, total=TOTAL,
        title="Карта функциональных блоков",
        image_path=dia(0),
        intro="10 функциональных блоков: F1–F5 — операционная цепочка, F6–F10 — экспертиза и сервисы, F7 — обеспечение.",
        bullets=[
            ("F1–F5 — ", "sourcing → QC → образцы → проблемы → liaison."),
            ("F6 — ", "технология детской одежды (tech pack, PPS, GB 31701)."),
            ("F8–F10 — ", "контент, лаборатории, аналитика рынка."),
            ("F7 — ", "администрация, учёт WFOE, биллинг материнской."),
        ],
    ); p += 1

    # --- Слайд 5. Операционная цепочка F1–F5 + KPI ---
    slide_table(
        prs, page=p, total=TOTAL,
        title="Операционная цепочка F1–F5 и KPI",
        headers=["Функция", "Цель", "KPI (целевые)"],
        rows=[
            ["F1. Sourcing", "Поддержка ПМ: фабрики, AVL, КП", "Ответ ПМ ≤ 2 дня; ≥ 90% брифов в срок"],
            ["F2. ОТК / QC", "Pre-shipment + аудиты до отгрузки", "Дефектность < 1,5%; PSI ≤ 1 день"],
            ["F3. Образцы", "Преселекция и отправка в РФ", "3–7 дней «фабрика → отправка»; ≥ 98%"],
            ["F4. Брак/задержки", "Выезд, root cause, CAPA", "Выезд ≤ 3 дня; ≥ 80% без эскалации"],
            ["F5. Liaison", "Единое окно для Москвы", "Ответ ≤ 1 день; статус без срывов"],
        ],
        col_widths=[1.5, 3.0, 3.5],
        note="Доп. функции (M5+): F6 tech pack/PPS · F8 контент к старту продаж · F9 тесты ТР ТС/GB · F10 тренды рынка.",
        body_size=12,
    ); p += 1

    # --- Слайд 6. Оргструктура ---
    slide_diagram(
        prs, page=p, total=TOTAL,
        title="Организационная структура (целевая, к M6)",
        image_path=dia(1),
        caption="Сплошные линии — административное подчинение директору; пунктир — функциональные связи с Москвой и внутри команды.",
    ); p += 1

    # --- Слайд 7. Численность и локации ---
    slide_table(
        prs, page=p, total=TOTAL,
        title="Численность и локации (к концу M6)",
        intro="Модель найма: RU — директор (Z-виза, legal representative) и технолог; CN — все операционные роли.",
        headers=["Локация", "Роли", "Чел."],
        rows=[
            ["Шанхай (HQ)", "Директор (RU), технолог (RU), бухгалтер, офис-админ, 3 кат. менеджера, контент, тестирование, аналитик", "10"],
            ["Гуанчжоу (QC-узел)", "Руководитель ОТК + инспекторы (1–2)", "2–3"],
            ["Итого", "", "12–13"],
        ],
        col_widths=[1.6, 5.4, 0.8],
        right_align_cols={2},
        emphasis_rows={2},
        note="Кат. менеджеры и технолог базируются в Шанхае, но часто в командировках (Гуандун, Фуцзянь, Чжэцзян, Чэнхай).",
        body_size=12,
    ); p += 1

    # --- Слайд 8. Бюджет — сводка ---
    slide_table(
        prs, page=p, total=TOTAL,
        title="Бюджет первых 6 месяцев — сводка",
        intro="Курс-ориентир: 1 CNY = 12 RUB · модель cost-plus 8%.",
        headers=["Блок", "CNY", "RUB (×12)"],
        rows=[
            ["CAPEX (единовременно, старт)", "254 000", "3 048 000"],
            ["Себестоимость услуг — cost base (6 мес.)", "1 730 675", "20 768 100"],
            ["Налоги WFOE (НДС 6% + CIT 25%)", "146 761", "1 761 132"],
            ["Выручка по сервисным договорам (без НДС)", "1 869 129", "22 429 548"],
            ["Run-rate на M6 (cost base / мес)", "433 025", "5 196 300"],
        ],
        col_widths=[4.2, 1.6, 1.8],
        right_align_cols={1, 2},
        emphasis_rows={3},
        note="Денежный отток материнской (6 мес.): выручка + НДС = 1 981 277 CNY (~23,78 млн RUB). CAPEX в M1–M2 покрывается инъекцией оборотного капитала.",
        body_size=12,
    ); p += 1

    # --- Слайд 9. Финансовая модель cost-plus ---
    slide_diagram(
        prs, page=p, total=TOTAL,
        title="Финансовая модель: сервисные договоры (cost-plus)",
        image_path=dia(2),
        intro="WFOE не получает «дотацию», а продаёт услуги материнской компании. Наценка 8% — ориентир для трансфертного ценообразования.",
        caption="Один рамочный сервисный договор + приложения-SOW по услугам + ежемесячные акты (биллинг в CNY).",
    ); p += 1

    # --- Слайд 10. Перечень сервисных договоров ---
    slide_table(
        prs, page=p, total=TOTAL,
        title="Перечень сервисных договоров (billable services)",
        headers=["Услуга", "Функция", "Модель оплаты"],
        rows=[
            ["Поиск и верификация поставщиков (sourcing fee)", "F1", "Ретейнер / за бриф"],
            ["Контроль качества и аудиты (QC / audit fee)", "F2, F4", "За инспекцию/аудит + ретейнер"],
            ["Управление и отправка образцов (sample fee)", "F3", "За партию / ретейнер"],
            ["Создание контента (content production fee)", "F8", "За SKU / пакет"],
            ["Координация тестирования (lab testing fee)", "F9", "Координация + pass-through"],
            ["Аналитика рынка (market intelligence retainer)", "F10", "Ежемесячный ретейнер"],
            ["Управление и liaison (management fee)", "F5, F7", "Ежемесячный ретейнер"],
        ],
        col_widths=[4.2, 1.2, 2.6],
        note="Счета лабораторий SGS/BV/Intertek — pass-through: перевыставляются отдельной строкой, в cost base не входят.",
        body_size=11.5,
    ); p += 1

    # --- Слайд 11. Roadmap M1–M6 (диаграмма + таблица) ---
    slide_diagram_table(
        prs, page=p, total=TOTAL,
        title="Roadmap запуска M1–M6",
        image_path=dia(3),
        headers=["Месяц", "Ключевые действия", "Веха"],
        rows=[
            ["M1", "Директор, юрконсультант, ТМ в CNIPA, проработка сервисных договоров", "Scope утверждён"],
            ["M2", "Pre-approval WFOE, офис-админ, кат. менеджер (игрушки), рамочный договор", "Команда 3 чел."],
            ["M3", "Business license, счёт, экспортная рег. «про запас», протокол ОТК AQL 2,5", "WFOE действует, 6 чел."],
            ["M4", "Бухгалтер, инспектор, аналитик; первые образцы и PSI; первый биллинг", "9 чел., первый акт"],
            ["M5", "Технолог, контент, тестирование, 2-й инспектор; PPS, контент, тесты", "13 чел., полная команда"],
            ["M6", "Штатный режим, KPI, ретроспектива, go/no-go экспорт (Фаза 2)", "KPI на уровне"],
        ],
        col_widths=[0.7, 5.0, 2.0],
        emphasis_rows={3, 4},
    ); p += 1

    # --- Слайд 12. Вехи, KPI, найм ---
    slide_two_column_tables(
        prs, page=p, total=TOTAL,
        title="KPI по фазам и график найма",
        left_title="KPI по фазам",
        left_headers=["KPI", "M1–2", "M3–4", "M5–6"],
        left_rows=[
            ["Численность", "1→3", "6→9", "13"],
            ["Фабрики (накопит.)", "0", "10–15", "25–40"],
            ["Образцы в РФ", "0", "первые", "поток"],
            ["PSI / аудиты в мес.", "0", "первые", "по SKU"],
            ["Дефектность отгрузок", "—", "базлайн", "< 1,5%"],
            ["SLA ответа ПМ", "—", "≤ 2 дня", "≤ 1 день"],
        ],
        left_cw=[1.6, 0.8, 0.8, 0.8],
        right_title="График найма (активных сотрудников)",
        right_headers=["Роль / месяц", "M1", "M2", "M3", "M4", "M5", "M6"],
        right_rows=[
            ["Директор (RU)", "●", "●", "●", "●", "●", "●"],
            ["Офис-админ (CN)", "", "●", "●", "●", "●", "●"],
            ["Кат. менеджер игрушки", "", "●", "●", "●", "●", "●"],
            ["Кат. менеджеры одежда/акс.", "", "", "●", "●", "●", "●"],
            ["Руководитель ОТК", "", "", "●", "●", "●", "●"],
            ["Бухгалтер, инспектор, аналитик", "", "", "", "●", "●", "●"],
            ["Технолог, контент, тестир., инсп.2", "", "", "", "", "●", "●"],
            ["Активных сотрудников", "1", "3", "6", "9", "13", "13"],
        ],
        right_cw=[2.2, 0.4, 0.4, 0.4, 0.4, 0.4, 0.4],
        note="13 вех M1–M6 — см. slides.md (слайд 12) и 04_setup_development_plan_6m.md. Поэтапный найм снижает burn rate первых месяцев.",
    ); p += 1

    # --- Слайд 13. Горизонтальные связи с Москвой ---
    slide_diagram(
        prs, page=p, total=TOTAL,
        title="Горизонтальные связи с московским офисом",
        image_path=dia(4),
        intro="Единое окно (директор WFOE, F5) + прямые функциональные связки кат. менеджеров и ПМ.",
        bullets=[
            ("Единое окно: ", "все запросы через директора — без «дёргания» исполнителей в обход приоритетов."),
            ("Функциональные связки: ", "кат. менеджеры ↔ ПМ; технолог ↔ ПМ одежды; контент ↔ маркетинг."),
            ("Общий трекер: ", "задачи, AVL, библиотека образцов, протоколы тестов — обеим сторонам. Часовые пояса: МСК +5 ч."),
        ],
    ); p += 1

    # --- Слайд 14. Сквозные процессы ---
    slide_three_processes(
        prs, page=p, total=TOTAL,
        title="Сквозные процессы: образцы, QC и оплата доставки",
        images=[dia(5), dia(6), dia(7)],
        captions=[
            "Процесс А. Образцы (3–7 дней «фабрика → РФ»)",
            "Процесс Б. Pre-shipment QC (AQL 2,5)",
            "Контур оплаты доставки образцов (cost-plus)",
        ],
        note="WFOE оплачивает перевозчика из своих средств → расход в cost base → перевыставление материнской с наценкой 8% + НДС.",
    ); p += 1

    # --- Слайд 15. SLA, риски, Фаза 2, финал ---
    slide_closing(
        prs, page=p, total=TOTAL,
        title="SLA, риски и критерии активации Фазы 2",
        headline="WFOE в Фазе 1 — контролируемый сервисный контур «глаза и руки» в Китае с прозрачным cost-plus финансированием, а не торговый риск с первого дня.",
        sla_rows=[
            ["Ответ на запрос ПМ", "≤ 1 день"],
            ["Закрытие брифа на образец", "3–7 дней"],
            ["Выезд по браку", "≤ 3 дня"],
            ["PSI после инспекции", "≤ 1 день"],
            ["Биллинг материнской", "ежемесячно"],
            ["Статус-отчёт", "еженедельно"],
            ["Аналитические обзоры", "раз в 2 недели"],
        ],
        risk_rows=[
            ["Задержка Z-визы", "Старт M1, HR-агентство"],
            ["Срыв регистрации WFOE", "Юрконсультант, буфер сроков"],
            ["Трансфертное ценообр.", "Согласовать cost-plus M1–M2"],
            ["Кража ИС", "ТМ в CNIPA до фабрик, NDA"],
            ["Открытие счёта (санкции)", "Bank of China / ICBC, CNY"],
        ],
        checklist=[
            "KPI сервисных функций на уровне, биллинг отлажен",
            "Подтверждённая экономия на закупках (Finance)",
            "Бизнес-кейс под торговые операции (B2B)",
            "Готовность Finance к товарным CNY-расчётам",
            "Положительная оценка санкционного/правового риска",
        ],
    )

    target = _save_pptx(prs, OUTPUT_PPTX)
    print(f"Сохранено: {target}")
    print(f"Всего слайдов: {len(prs.slides._sldIdLst)}")


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Сборка PPTX дека WFOE Детский Мир в КНР (Фаза 1).")
    parser.add_argument("--skip-mermaid", action="store_true",
                        help="Не рендерить mermaid, использовать готовые PNG из _assets/diagrams/")
    args = parser.parse_args()

    diagrams = render_mermaid(skip=args.skip_mermaid)
    build_deck(diagrams)


if __name__ == "__main__":
    main()
